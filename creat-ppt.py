from io import BytesIO
import streamlit as st
from PIL import Image
import uuid
import hashlib
import tempfile
import os
import time
import imagehash

# ================= 页面配置 =================
st.set_page_config(layout="wide", page_title="多图片自动生成 PPT")
st.title("多图片自动生成 PPT")

# ================= 初始化 Session State =================
def init_session():
    if "images" not in st.session_state:
        st.session_state.images = []
    if "processed_ids" not in st.session_state:
        st.session_state.processed_ids = set()
    if "page" not in st.session_state:
        st.session_state.page = 1
    if "ppt_bytes" not in st.session_state:
        st.session_state.ppt_bytes = None
    if "temp_duplicates" not in st.session_state:
        st.session_state.temp_duplicates = []
    # 通过 key 的变动来强制清空 file_uploader
    if "uploader_key" not in st.session_state:
        st.session_state.uploader_key = str(uuid.uuid4())

init_session()

TMP_DIR = tempfile.gettempdir()

# ================= 核心功能：清空功能 =================
def clear_all_data():
    # 物理删除临时缩略图
    for img in st.session_state.images:
        try:
            if os.path.exists(img["thumb_path"]):
                os.remove(img["thumb_path"])
        except:
            pass
    
    # 重置状态
    st.session_state.images = []
    st.session_state.processed_ids = set()
    st.session_state.page = 1
    st.session_state.ppt_bytes = None
    st.session_state.temp_duplicates = []
    # 核心：改变 key，彻底清空上传组件的文件列表
    st.session_state.uploader_key = str(uuid.uuid4())
    st.rerun()

# ================= 核心功能：重复检测弹窗 (修正版) =================
@st.dialog("发现疑似重复图片")
def show_duplicate_dialog():
    st.info("系统检测到图片内容相似。**已自动勾选分辨率较低（或重复）的版本**，您可以手动调整：")
    
    # 记录用户想要删除的 UID
    uids_to_remove = set()
    
    for idx, dup in enumerate(st.session_state.temp_duplicates):
        orig = dup['original']
        curr = dup['current']
        
        # --- 获取分辨率数据 ---
        orig_px = orig.get('pixels', 0)
        curr_px = curr.get('pixels', 0)
        orig_res_str = orig.get('res_str', '未知')
        curr_res_str = curr.get('res_str', '未知')

        # --- 自动勾选逻辑 (修正了相等情况) ---
        # 1. 如果旧图比新图糊 -> 删旧图
        # 2. 否则（新图糊，或者一样清晰）-> 删新图（保留原有）
        if orig_px < curr_px:
            default_del_orig = True
            default_del_curr = False
        else:
            default_del_orig = False
            default_del_curr = True

        col1, col2, col3 = st.columns([4, 1, 4])
        
        with col1:
            st.image(orig['thumb_path'], width=180)
            st.markdown(f"**已有图片**: {orig['name']}")
            st.caption(f"📏 分辨率: {orig_res_str}")
            
            # 删除已有
            if st.checkbox(f"删除这张 (已有)", value=default_del_orig, key=f"del_orig_{idx}_{orig['uid']}"):
                uids_to_remove.add(orig['uid'])
        
        with col2:
            st.markdown("<br><br><h3 style='text-align: center; color: gray;'>VS</h3>", unsafe_allow_html=True)
        
        with col3:
            st.image(curr['thumb_path'], width=180)
            st.markdown(f"**新上传项**: {curr['name']}")
            st.caption(f"📏 分辨率: {curr_res_str}")
            
            # 删除新传
            if st.checkbox(f"删除这张 (新传)", value=default_del_curr, key=f"del_curr_{idx}_{curr['uid']}"):
                uids_to_remove.add(curr['uid'])
        
        st.divider()

    if st.button("确认处理并关闭弹窗", type="primary", use_container_width=True):
        if uids_to_remove:
            st.session_state.images = [
                img for img in st.session_state.images 
                if img['uid'] not in uids_to_remove
            ]
        # 处理完后清空临时队列
        st.session_state.temp_duplicates = []
        st.success(f"已处理！成功删除 {len(uids_to_remove)} 张图片")
        time.sleep(0.5)
        st.rerun()

# ================= 顶部操作栏 =================
col_upload, col_clear = st.columns([8, 2])

with col_upload:
    uploaded_files = st.file_uploader(
        "上传图片（支持批量，Ctrl+A 全选）",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key=st.session_state.uploader_key
    )

with col_clear:
    st.write("---") 
    if st.button("🗑️ 清空所有图片", use_container_width=True, type="secondary"):
        clear_all_data()

# ================= 上传处理逻辑 (修改版) =================
if uploaded_files:
    if not st.session_state.temp_duplicates:
        new_found_duplicates = []
        SIMILARITY_THRESHOLD = 1 # 稍微放宽一点阈值，以免漏掉颜色有细微差别的重复项

        for file in uploaded_files:
            file_bytes = file.read()
            file_hash = hashlib.md5(file_bytes).hexdigest()
            file_id = f"{file.name}_{file_hash}"

            if file_id in st.session_state.processed_ids:
                continue

            try:
                img = Image.open(BytesIO(file_bytes)).convert("RGB")
                
                # --- 获取图片详细尺寸信息 ---
                width, height = img.size
                pixels = width * height
                res_str = f"{width} x {height}"
                
                curr_phash = imagehash.phash(img)
                
                uid = str(uuid.uuid4())
                thumb = img.copy()
                thumb.thumbnail((260, 260))
                thumb_path = os.path.join(TMP_DIR, f"{uid}.png")
                thumb.save(thumb_path, "PNG")

                new_img_obj = {
                    "uid": uid,
                    "name": file.name,
                    "bytes": file_bytes,
                    "thumb_path": thumb_path,
                    "ratio": width / height,
                    "phash": curr_phash,
                    # 新增字段：用于对比清晰度
                    "pixels": pixels,
                    "res_str": res_str
                }

                # 查找重复
                is_duplicate = False
                for existing in st.session_state.images:
                    if (curr_phash - existing['phash']) <= SIMILARITY_THRESHOLD:
                        new_found_duplicates.append({"original": existing, "current": new_img_obj})
                        is_duplicate = True
                        break 
                
                # 无论是否重复，先加入列表（如果是重复的，会在弹窗里决定是否删除）
                # 逻辑说明：如果不加入，后续处理会很麻烦。不如先加进去，然后通过弹窗删掉 UID。
                # 但这里的逻辑稍微调整一下：只把新图片存入 images，如果用户在弹窗选了删除新图片，再从 images 移除
                st.session_state.images.append(new_img_obj)
                st.session_state.processed_ids.add(file_id)

            except Exception as e:
                st.error(f"{file.name} 读取失败：{e}")

        if new_found_duplicates:
            st.session_state.temp_duplicates = new_found_duplicates
            show_duplicate_dialog()

# ================= 展示与分页 =================
IMAGES_PER_PAGE = 40
IMAGES_PER_ROW = 10
THUMB_HEIGHT_MM = 40
MM_TO_PIXELS = 3.77953

total_images = len(st.session_state.images)
total_pages = max(1, (total_images + IMAGES_PER_PAGE - 1) // IMAGES_PER_PAGE)
st.session_state.page = min(st.session_state.page, total_pages)

start_idx = (st.session_state.page - 1) * IMAGES_PER_PAGE
page_images = st.session_state.images[start_idx : start_idx + IMAGES_PER_PAGE]

st.subheader(f"图片预览 (共 {total_images} 张)")

# 网格展示
if len(page_images) > 0:
    for i in range(0, len(page_images), IMAGES_PER_ROW):
        cols = st.columns(IMAGES_PER_ROW)
        for col_idx, img_obj in enumerate(page_images[i:i + IMAGES_PER_ROW]):
            with cols[col_idx]:
                h_px = int(THUMB_HEIGHT_MM * MM_TO_PIXELS)
                w_px = int(h_px * img_obj["ratio"])
                st.image(img_obj["thumb_path"], width=w_px)
else:
    st.info("暂无图片，请上传。")

# 分页导航
if total_pages > 1:
    cp, cn, _ = st.columns([1, 1, 6])
    with cp:
        if st.button("上一页", disabled=(st.session_state.page <= 1)):
            st.session_state.page -= 1
            st.rerun()
    with cn:
        if st.button("下一页", disabled=(st.session_state.page >= total_pages)):
            st.session_state.page += 1
            st.rerun()

# ================= PPT 生成 =================
def generate_ppt(images):
    from pptx import Presentation
    from pptx.util import Inches, Mm
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    left_m, top_m, space, fix_h = Mm(0), Mm(10), Mm(2.5), Mm(40)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    x, y = left_m, top_m
    
    for img in images:
        w = fix_h * img["ratio"]
        if x + w > prs.slide_width:
            x, y = left_m, y + fix_h + space
        if y + fix_h > prs.slide_height - top_m:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            x, y = left_m, top_m
        
        slide.shapes.add_picture(BytesIO(img["bytes"]), x, y, height=fix_h)
        x += w + space
    
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

st.divider()
if st.button("🚀 生成 PPT", type="primary", use_container_width=True):
    if st.session_state.images:
        with st.spinner("正在排版生成中..."):
            st.session_state.ppt_bytes = generate_ppt(st.session_state.images)
        st.success("PPT 生成成功！")
    else:
        st.warning("请先上传图片")

if st.session_state.ppt_bytes:
    st.download_button(
        "📂 下载 PPT 文件", 
        data=st.session_state.ppt_bytes, 
        file_name=f"ppt_export_{int(time.time())}.pptx", 
        use_container_width=True
    )
